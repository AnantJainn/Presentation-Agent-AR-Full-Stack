# agents/pptx_agent.py
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
import requests
from io import BytesIO
import urllib.parse

def get_ai_image(prompt):
    """
    Generates an image on the fly using Pollinations.ai (Free, No API Key).
    This ensures the image ALWAYS matches the specific domain/topic.
    """
    clean_prompt = urllib.parse.quote(prompt)
    url = f"https://image.pollinations.ai/prompt/{clean_prompt}?width=1024&height=1024&nologo=true"
    
    try:
        # User-Agent headers are critical to avoid being blocked
        headers = {'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36'}
        response = requests.get(url, headers=headers, timeout=10)
        response.raise_for_status()
        return BytesIO(response.content)
    except Exception as e:
        print(f"⚠️ Image gen failed for '{prompt}': {e}")
        return None

def pptx_agent(state):
    print("🎨 Generative PPTX Engine Starting...")
    prs = Presentation()
    # Widescreen 16:9
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    slides_data = state["presentation"]["slides"]

    for i, slide_data in enumerate(slides_data):
        slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank layout
        
        # --- 1. DESIGN BACKGROUND ---
        # Main background (Soft Off-White)
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = RGBColor(250, 250, 252)

        # Accent Shape (Top Left organic blob or bar)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 
            Inches(0), Inches(0), Inches(13.33), Inches(0.15)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(50, 50, 200) # Professional Blue
        shape.line.fill.background()

        # --- 2. TITLE (Fixed Overflow) ---
        # We give the title the full width of the slide minus margins
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5), Inches(12.33), Inches(1.2)
        )
        tf = title_box.text_frame
        tf.word_wrap = True # CRITICAL: Wraps text if it's too long
        
        p = tf.paragraphs[0]
        p.text = slide_data["title"]
        p.font.size = Pt(36) # Large readable font
        p.font.bold = True
        p.font.name = "Arial"
        p.font.color.rgb = RGBColor(30, 30, 30)

        # --- 3. IMAGE (Generated on the fly) ---
        # Placed on the Right side
        visual_prompt = slide_data.get("visual_description", "abstract technology background")
        print(f"   🖼️ Generating image: {visual_prompt}...")
        image_stream = get_ai_image(visual_prompt)
        
        if image_stream:
            # Add image (Right side, square-ish)
            pic = slide.shapes.add_picture(
                image_stream, 
                Inches(8.0), Inches(2.0),
                width=Inches(4.8), height=Inches(4.8)
            )
            # Add a subtle border to the image
            line = pic.line
            line.color.rgb = RGBColor(200, 200, 200)
            line.width = Pt(1)

        # --- 4. CONTENT (Left side) ---
        # Main Paragraph
        content_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.8), Inches(7.0), Inches(3.0)
        )
        tf = content_box.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = slide_data.get("content", "")
        p.font.size = Pt(18)
        p.font.name = "Georgia" # Good for body text
        p.line_spacing = 1.3

        # Key Points (Bottom Left)
        if "key_points" in slide_data:
            # Subtle background for key points
            kp_bg = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.5), Inches(5.0), Inches(7.0), Inches(2.0)
            )
            kp_bg.fill.solid()
            kp_bg.fill.fore_color.rgb = RGBColor(240, 244, 255) # Light Blue tint
            kp_bg.line.fill.background()

            kp_box = slide.shapes.add_textbox(
                Inches(0.6), Inches(5.1), Inches(6.8), Inches(1.8)
            )
            tf = kp_box.text_frame
            tf.word_wrap = True
            
            for point in slide_data["key_points"]:
                p = tf.add_paragraph()
                p.text = f"• {point}"
                p.font.size = Pt(14)
                p.font.name = "Arial"
                p.space_after = Pt(10)

    # Save
    output_file = "output.pptx"
    prs.save(output_file)
    print(f"✅ Generated {output_file} with {len(slides_data)} slides.")
    return state